from server import mcp
import os
from utils.create_files import create_file
from utils.create_folders import create_folder

@mcp.tool()
def create_txt_file_tool(file_name: str = "", content: str = ""):
    default_path = r"C:\Users\emirc\OneDrive\Masaüstü"
    if not file_name:
        file_name = "yeni_dosya.txt"
    else:
        file_name = f"{file_name}.txt"
    path = os.path.join(default_path, file_name)
    
    # İçerik varsa dosyaya yaz, yoksa varsayılan içerik
    if content:
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
    else:
        create_file(path)
    return f"TXT dosyası oluşturuldu: {path}"

@mcp.tool()
def create_word_file_tool(file_name: str = "", content: str = ""):
    default_path = r"C:\Users\emirc\OneDrive\Masaüstü"
    if not file_name:
        file_name = "yeni_dosya.docx"
    else:
        file_name = f"{file_name}.docx"
    path = os.path.join(default_path, file_name)
    
    # Word dosyası oluştur ve içerik ekle
    from docx import Document
    doc = Document()
    
    if content:
        # İçeriği paragraflara böl ve ekle
        paragraphs = content.split('\n')
        for i, para in enumerate(paragraphs):
            if i == 0 and para.strip():  # İlk paragraf başlık olsun
                doc.add_heading(para.strip(), 0)
            elif para.strip():
                doc.add_paragraph(para.strip())
    else:
        doc.add_heading('Başlık', 0)
        doc.add_paragraph('Bu bir Word dosyasıdır.')
    
    doc.save(path)
    return f"Word dosyası oluşturuldu: {path}"

@mcp.tool()
def create_excel_file_tool(file_name: str = "", content: str = ""):
    try:
        default_path = r"C:\Users\emirc\OneDrive\Masaüstü"
        if not file_name:
            file_name = "yeni_dosya.xlsx"
        else:
            file_name = f"{file_name}.xlsx"
        path = os.path.join(default_path, file_name)
        
        # Excel dosyası oluştur ve içerik ekle
        try:
            from openpyxl import Workbook
        except ImportError:
            return "Hata: openpyxl kütüphanesi yüklü değil. Lütfen 'pip install openpyxl' komutunu çalıştırın."
        
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Veri"
        
        if content:
            # İçeriği satırlara böl ve ekle
            rows = content.split('\n')
            for i, row in enumerate(rows, 1):
                if row.strip():
                    # Farklı ayırıcı karakterleri kontrol et
                    if '\t' in row:
                        # Tab ile ayrılmış veriler
                        cells = row.split('\t')
                    elif ';' in row:
                        # Noktalı virgül ile ayrılmış veriler
                        cells = row.split(';')
                    elif ',' in row:
                        # Virgül ile ayrılmış veriler (CSV formatı)
                        import csv
                        from io import StringIO
                        
                        # StringIO kullanarak CSV reader ile parse et
                        csv_reader = csv.reader(StringIO(row))
                        cells = next(csv_reader)  # İlk (ve tek) satırı al
                    elif '  ' in row or '\t' in row:
                        # Birden fazla boşluk veya tab ile ayrılmış veriler
                        import re
                        cells = re.split(r'\s{2,}|\t', row)
                    else:
                        # Tek boşluk ile ayrılmış veriler
                        cells = row.split(' ')
                    
                    # Her hücreyi ilgili sütuna yaz
                    for j, cell in enumerate(cells, 1):
                        sheet.cell(row=i, column=j, value=cell.strip())
        else:
            sheet['A1'] = "Yeni Excel Dosyası"
        
        workbook.save(path)
        return f"Excel dosyası başarıyla oluşturuldu: {path}"
    
    except Exception as e:
        return f"Excel dosyası oluşturulurken hata oluştu: {str(e)}"

@mcp.tool()
def create_powerpoint_file_tool(file_name: str = "", content: str = ""):
    default_path = r"C:\Users\emirc\OneDrive\Masaüstü"
    if not file_name:
        file_name = "yeni_dosya.pptx"
    else:
        file_name = f"{file_name}.pptx"
    path = os.path.join(default_path, file_name)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    prs = Presentation()
    if content:
        import re
        
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        slides = []
        
        # Slayt algılama desenleri
        slide_patterns = [
            re.compile(r'^(slayt|slide)\s*\d+.*', re.IGNORECASE),  # Slayt 1, Slide 2, vb.
            re.compile(r'^\d+\.\s+.*', re.IGNORECASE),  # 1. Başlık, 2. Başlık, vb.
            re.compile(r'^[A-Z][A-Z\s]+$', re.IGNORECASE),  # SADECE BÜYÜK HARFLERLE YAZILMIŞ BAŞLIKLAR
            re.compile(r'^[A-Z][a-z\s]+:$', re.IGNORECASE),  # Başlık: formatı
        ]
        
        current_title = None
        current_body = []
        
        for line in lines:
            is_title = False
            
            # Slayt numarası ile başlayan satırları kontrol et
            for pattern in slide_patterns:
                if pattern.match(line):
                    if current_title is not None:
                        slides.append((current_title, '\n'.join(current_body).strip()))
                    current_title = line
                    current_body = []
                    is_title = True
                    break
            
            # Eğer başlık değilse, içeriğe ekle
            if not is_title:
                current_body.append(line)
        
        # Son slaytı ekle
        if current_title is not None:
            slides.append((current_title, '\n'.join(current_body).strip()))
        
        # Eğer hiç slayt bulunamadıysa, içeriği paragraflara böl
        if not slides:
            # İçeriği paragraflara böl ve her paragrafı ayrı slayt yap
            paragraphs = [p for p in lines if p.strip()]
            if len(paragraphs) > 1:
                # İlk paragraf başlık olsun
                slides.append((paragraphs[0], '\n'.join(paragraphs[1:])))
            else:
                slides = [("Başlık", content)]
        
        MAX_LINES_PER_SLIDE = 15
        for title_text, body_text in slides:
            # Başlık metnini temizle
            clean_title = title_text
            
            # Slayt numarası varsa temizle
            slide_match = re.match(r'^(slayt|slide)\s*\d+\s*[:\-]?\s*(.*)', title_text, re.IGNORECASE)
            if slide_match:
                clean_title = slide_match.group(2) if slide_match.group(2) else slide_match.group(0)
                if not clean_title.strip():
                    clean_title = slide_match.group(0)
            
            # Numaralı başlık formatını temizle (1. Başlık -> Başlık)
            num_match = re.match(r'^\d+\.\s+(.*)', clean_title)
            if num_match:
                clean_title = num_match.group(1)
            
            # Sondaki iki noktayı temizle
            clean_title = clean_title.rstrip(':')
            
            # Eğer başlık çok uzunsa kısalt
            if len(clean_title) > 50:
                clean_title = clean_title[:47] + "..."
            
            paragraphs = [p for p in body_text.split('\n') if p.strip()]
            
            # Eğer çok fazla içerik varsa birden fazla slayta böl
            for i in range(0, len(paragraphs), MAX_LINES_PER_SLIDE):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Boş slayt
                
                # Başlık ekle
                left = Inches(0.5)
                top = Inches(0.6)
                width = Inches(9)
                height = Inches(1)
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_tf = title_box.text_frame
                title_tf.text = clean_title
                title_tf.paragraphs[0].font.size = Pt(40)
                title_tf.paragraphs[0].font.bold = True
                title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # İçerik ekle
                content_top = Inches(2.0)
                content_height = Inches(4.5)
                content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
                content_tf = content_box.text_frame
                content_tf.word_wrap = True
                content_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                
                for para in paragraphs[i:i+MAX_LINES_PER_SLIDE]:
                    p = content_tf.add_paragraph()
                    p.text = para
                    p.font.size = Pt(20)
                    p.alignment = PP_ALIGN.LEFT
                
                # İlk paragraf boş gelir, onu sil
                if content_tf.paragraphs and not content_tf.paragraphs[0].text:
                    content_tf._element.remove(content_tf.paragraphs[0]._element)
        
        prs.save(path)
        return f"PowerPoint dosyası oluşturuldu: {path}"
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "PowerPoint Başlığı"
    prs.save(path)
    return f"PowerPoint dosyası oluşturuldu: {path}"

@mcp.tool()
def create_folder_tool(folder_name: str = "", content: str = "", file_name: str = "", file_type: str = "txt"):
    default_path = r"C:\Users\emirc\OneDrive\Masaüstü"
    if not folder_name:
        folder_name = "YeniKlasor"
    
    # Klasör yolunu oluştur
    folder_path = os.path.join(default_path, folder_name)
    create_folder(folder_path)
    
    result_message = f"Klasör oluşturuldu: {folder_path}"
    
    # Eğer dosya adı ve içerik verilmişse, klasör içine dosya oluştur
    if file_name and content:
        # Dosya tipine göre uzantı belirle
        if file_type.lower() == "word" or file_type.lower() == "docx":
            file_extension = ".docx"
        elif file_type.lower() == "excel" or file_type.lower() == "xlsx":
            file_extension = ".xlsx"
        elif file_type.lower() == "powerpoint" or file_type.lower() == "pptx":
            file_extension = ".pptx"
        else:
            file_extension = ".txt"
        
        # Dosya adına uzantı ekle (eğer yoksa)
        if not file_name.endswith(file_extension):
            file_name = f"{file_name}{file_extension}"
        
        # Dosya yolunu oluştur
        file_path = os.path.join(folder_path, file_name)
        
        # Dosya tipine göre oluştur
        if file_type.lower() in ["word", "docx"]:
            from docx import Document
            doc = Document()
            paragraphs = content.split('\n')
            for i, para in enumerate(paragraphs):
                if i == 0 and para.strip():
                    doc.add_heading(para.strip(), 0)
                elif para.strip():
                    doc.add_paragraph(para.strip())
            doc.save(file_path)
            result_message += f"\nWord dosyası oluşturuldu: {file_path}"
            
        elif file_type.lower() in ["excel", "xlsx"]:
            from openpyxl import Workbook
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Veri"
            rows = content.split('\n')
            for i, row in enumerate(rows, 1):
                if row.strip():
                    cells = row.split(',')
                    for j, cell in enumerate(cells, 1):
                        sheet.cell(row=i, column=j, value=cell.strip())
            workbook.save(file_path)
            result_message += f"\nExcel dosyası oluşturuldu: {file_path}"
            
        elif file_type.lower() in ["powerpoint", "pptx"]:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
            prs = Presentation()
            import re
            
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            slides = []
            
            # Slayt algılama desenleri
            slide_patterns = [
                re.compile(r'^(slayt|slide)\s*\d+.*', re.IGNORECASE),  # Slayt 1, Slide 2, vb.
                #re.compile(r'^\d+\.\s+.*', re.IGNORECASE),  # 1. Başlık, 2. Başlık, vb.
                re.compile(r'^[A-Z][A-Z\s]+$', re.IGNORECASE),  # SADECE BÜYÜK HARFLERLE YAZILMIŞ BAŞLIKLAR
            ]
            
            current_title = None
            current_body = []
            
            for line in lines:
                is_title = False
                
                # Slayt numarası ile başlayan satırları kontrol et
                for pattern in slide_patterns:
                    if pattern.match(line):
                        if current_title is not None:
                            slides.append((current_title, '\n'.join(current_body).strip()))
                        current_title = line
                        current_body = []
                        is_title = True
                        break
                
                # Eğer başlık değilse, içeriğe ekle
                if not is_title:
                    current_body.append(line)
            
            # Son slaytı ekle
            if current_title is not None:
                slides.append((current_title, '\n'.join(current_body).strip()))
            
            # Eğer hiç slayt bulunamadıysa, içeriği paragraflara böl
            if not slides:
                # İçeriği paragraflara böl ve her paragrafı ayrı slayt yap
                paragraphs = [p for p in lines if p.strip()]
                if len(paragraphs) > 1:
                    # İlk paragraf başlık olsun
                    slides.append((paragraphs[0], '\n'.join(paragraphs[1:])))
                else:
                    slides = [("Başlık", content)]
            
            MAX_LINES_PER_SLIDE = 15
            for title_text, body_text in slides:
                # Başlık metnini temizle
                clean_title = title_text
                
                # Slayt numarası varsa temizle
                slide_match = re.match(r'^(slayt|slide)\s*\d+\s*[:\-]?\s*(.*)', title_text, re.IGNORECASE)
                if slide_match:
                    clean_title = slide_match.group(2) if slide_match.group(2) else slide_match.group(0)
                    if not clean_title.strip():
                        clean_title = slide_match.group(0)
                
                # Numaralı başlık formatını temizle (1. Başlık -> Başlık)
                num_match = re.match(r'^\d+\.\s+(.*)', clean_title)
                if num_match:
                    clean_title = num_match.group(1)
                
                # Sondaki iki noktayı temizle
                clean_title = clean_title.rstrip(':')
                
                # Eğer başlık çok uzunsa kısalt
                if len(clean_title) > 50:
                    clean_title = clean_title[:47] + "..."
                
                paragraphs = [p for p in body_text.split('\n') if p.strip()]
                
                # Eğer çok fazla içerik varsa birden fazla slayta böl
                for i in range(0, len(paragraphs), MAX_LINES_PER_SLIDE):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Boş slayt
                    
                    # Başlık ekle
                    left = Inches(0.5)
                    top = Inches(0.6)
                    width = Inches(9)
                    height = Inches(1)
                    title_box = slide.shapes.add_textbox(left, top, width, height)
                    title_tf = title_box.text_frame
                    title_tf.text = clean_title
                    title_tf.paragraphs[0].font.size = Pt(40)
                    title_tf.paragraphs[0].font.bold = True
                    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    
                    # İçerik ekle
                    content_top = Inches(2.0)
                    content_height = Inches(4.5)
                    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
                    content_tf = content_box.text_frame
                    content_tf.word_wrap = True
                    content_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    
                    for para in paragraphs[i:i+MAX_LINES_PER_SLIDE]:
                        p = content_tf.add_paragraph()
                        p.text = para
                        p.font.size = Pt(20)
                        p.alignment = PP_ALIGN.LEFT
                    
                    # İlk paragraf boş gelir, onu sil
                    if content_tf.paragraphs and not content_tf.paragraphs[0].text:
                        content_tf._element.remove(content_tf.paragraphs[0]._element)
            
            prs.save(file_path)
            result_message += f"\nPowerPoint dosyası oluşturuldu: {file_path}"
            
        else:  # txt dosyası
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(content)
            result_message += f"\nTXT dosyası oluşturuldu: {file_path}"
    
    return result_message
        